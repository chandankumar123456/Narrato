from pptx import Presentation # type: ignore
from pptx.util import Inches, Pt, Emu # type: ignore
from pptx.dml.color import RGBColor # type: ignore
from pptx.enum.text import PP_ALIGN # type: ignore
import uuid, os, logging
from config import settings
from ppt.themes.modern import THEMES, ThemeConfig

logger = logging.getLogger(__name__)

SLIDE_W = Inches(13.33)  # Widescreen 16:9
SLIDE_H = Inches(7.5)

def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_text_box(slide, text, left, top, width, height, font_name, font_size, bold=False, color="000000", align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return txBox

def set_background(slide, color_hex: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)

def _inject_speaker_notes(slide, notes_text: str):
    """Inject speaker notes into a slide's notes page."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

def generate_ppt(state) -> str:
    theme: ThemeConfig = THEMES.get(state.theme, THEMES["modern"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    # Build a lookup for speaker notes by slide_id
    notes_map = {}
    if state.speaker_notes:
        for entry in state.speaker_notes:
            notes_map[entry.get("slide_id")] = entry.get("notes", "")

    for slide_data in state.structured_slides:
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide, theme.background)
        _render_slide(slide, slide_data, theme)

        # Inject speaker notes
        slide_id = slide_data.get("slide_id")
        notes_text = notes_map.get(slide_id, "")
        if notes_text:
            try:
                _inject_speaker_notes(slide, notes_text)
            except Exception:
                logger.warning("Failed to inject speaker notes for slide %s", slide_id)

    os.makedirs(settings.output_dir, exist_ok=True)
    path = f"{settings.output_dir}/{uuid.uuid4().hex}.pptx"
    prs.save(path)
    return path

def _render_slide(slide, slide_data: dict, theme: ThemeConfig):
    from ppt.layouts import (
        title_slide, section_header, agenda_slide, problem_slide,
        stats_slide, feature_slide, comparison_slide, timeline_slide,
        example_slide, quote_slide, image_slide, conclusion_slide,
        cta_slide, thank_you_slide,
    )

    dispatch = {
        "title_slide": title_slide.render,
        "section_header": section_header.render,
        "agenda_slide": agenda_slide.render,
        "problem_slide": problem_slide.render,
        "stats_slide": stats_slide.render,
        "feature_slide": feature_slide.render,
        "comparison_slide": comparison_slide.render,
        "timeline_slide": timeline_slide.render,
        "example_slide": example_slide.render,
        "quote_slide": quote_slide.render,
        "image_slide": image_slide.render,
        "conclusion_slide": conclusion_slide.render,
        "cta_slide": cta_slide.render,
        "thank_you_slide": thank_you_slide.render,
    }

    renderer = dispatch.get(slide_data.get("type"))
    if renderer:
        try:
            renderer(slide, slide_data.get("content", {}), theme, slide_data.get("image_path"))
        except Exception:
            logger.exception("Failed to render slide type %s, using fallback", slide_data.get("type"))
            _render_fallback(slide, slide_data, theme)
    else:
        _render_fallback(slide, slide_data, theme)

def _render_fallback(slide, slide_data: dict, theme: ThemeConfig):
    """Generic fallback layout for unknown or failed slide types."""
    content = slide_data.get("content", {})
    title = content.get("title", "") or content.get("section_title", "") or ""
    body = content.get("body", "") or content.get("description", "") or ""
    add_text_box(slide, title, Inches(0.5), Inches(0.3), Inches(12), Inches(1),
                 theme.font_heading, theme.heading_size, bold=True, color=theme.primary)
    if body:
        add_text_box(slide, body, Inches(0.5), Inches(1.5), Inches(12), Inches(5),
                     theme.font_body, theme.body_size, color=theme.text)