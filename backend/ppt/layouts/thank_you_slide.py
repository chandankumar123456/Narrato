from pptx.util import Inches, Pt  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from ppt.generator import add_text_box, hex_to_rgb, set_background


def render(slide, content: dict, theme, image_path=None):
    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme.primary)
    bar.line.fill.background()

    # Title
    add_text_box(
        slide, content.get("title", "Thank You"),
        Inches(1), Inches(2.2), Inches(11.33), Inches(1.5),
        theme.font_heading, 48, bold=True,
        color=theme.primary, align=PP_ALIGN.CENTER,
    )

    # Thin accent line
    line = slide.shapes.add_shape(
        1, Inches(5.67), Inches(3.8), Inches(2), Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(theme.accent)
    line.line.fill.background()

    # Message
    add_text_box(
        slide, content.get("message", ""),
        Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.0),
        theme.font_body, 22, bold=False,
        color=theme.secondary, align=PP_ALIGN.CENTER,
    )

    # Contact
    add_text_box(
        slide, content.get("contact", ""),
        Inches(1.5), Inches(5.5), Inches(10.33), Inches(0.6),
        theme.font_body, 16, bold=False,
        color=theme.text, align=PP_ALIGN.CENTER,
    )

    # Bottom accent bar
    bar2 = slide.shapes.add_shape(1, Inches(0), Inches(7.35), Inches(13.33), Inches(0.15))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = hex_to_rgb(theme.accent)
    bar2.line.fill.background()
