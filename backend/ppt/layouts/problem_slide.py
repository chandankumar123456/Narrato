from pptx.util import Inches, Pt  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from ppt.generator import add_text_box, hex_to_rgb, set_background


def render(slide, content: dict, theme, image_path=None):
    # Title
    add_text_box(
        slide, content.get("title", ""),
        Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.9),
        theme.font_heading, theme.heading_size, bold=True,
        color=theme.primary, align=PP_ALIGN.LEFT,
    )

    # Accent underline
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(1.25), Inches(1.5), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme.accent)
    bar.line.fill.background()

    # Up to 3 cards
    cards = content.get("cards", [])[:3]
    card_w = 3.5
    gap = 0.5
    count = len(cards) if cards else 1
    total_w = count * card_w + (count - 1) * gap
    start_x = (13.33 - total_w) / 2

    for i, card in enumerate(cards):
        x = start_x + i * (card_w + gap)
        y = 2.0

        # Card background
        bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w), Inches(4.8))
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(theme.background)
        bg.line.color.rgb = hex_to_rgb(theme.secondary)
        bg.line.width = Pt(1)

        # Top accent strip on card
        strip = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(card_w), Inches(0.1),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex_to_rgb(theme.accent)
        strip.line.fill.background()

        # Icon
        add_text_box(
            slide, card.get("icon", "●"),
            Inches(x), Inches(y + 0.5), Inches(card_w), Inches(0.7),
            theme.font_body, 28, bold=False,
            color=theme.accent, align=PP_ALIGN.CENTER,
        )

        # Label
        add_text_box(
            slide, card.get("label", ""),
            Inches(x + 0.25), Inches(y + 1.3), Inches(card_w - 0.5), Inches(0.7),
            theme.font_heading, 18, bold=True,
            color=theme.primary, align=PP_ALIGN.CENTER,
        )

        # Description
        add_text_box(
            slide, card.get("description", ""),
            Inches(x + 0.25), Inches(y + 2.1), Inches(card_w - 0.5), Inches(2.4),
            theme.font_body, 14, bold=False,
            color=theme.text, align=PP_ALIGN.CENTER,
        )
