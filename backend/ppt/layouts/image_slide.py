import os

from pptx.util import Inches  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore

from ppt.generator import add_text_box, hex_to_rgb
from ppt.design_system import Grid, Spacing, Typography, SLIDE_WIDTH, SLIDE_HEIGHT


def render(slide, content: dict, theme, image_path=None):
    # Full-bleed image
    if image_path and os.path.isfile(image_path):
        slide.shapes.add_picture(
            image_path, Inches(0), Inches(0),
            Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT),
        )

    # Overlay bar at bottom
    overlay_y = 5.0
    overlay_h = SLIDE_HEIGHT - overlay_y
    overlay = slide.shapes.add_shape(
        1, Inches(0), Inches(overlay_y),
        Inches(SLIDE_WIDTH), Inches(overlay_h),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = hex_to_rgb(theme.primary)
    overlay.line.fill.background()

    # Title on overlay
    left, width = Grid.full_width()
    add_text_box(
        slide, content.get("title", ""),
        Inches(left), Inches(overlay_y + Spacing.TIGHT),
        Inches(width), Inches(1.0),
        theme.font_heading, Typography.SUBHEADING + 2, bold=True,
        color=theme.background, align=PP_ALIGN.LEFT,
    )

    # Caption on overlay
    add_text_box(
        slide, content.get("caption", ""),
        Inches(left), Inches(overlay_y + 1.0 + Spacing.ELEMENT),
        Inches(width), Inches(0.8),
        theme.font_body, Typography.BODY, bold=False,
        color=theme.background, align=PP_ALIGN.LEFT,
    )
