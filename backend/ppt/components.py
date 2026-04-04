"""
Reusable visual components for slide rendering.

Each component follows the design system's spacing, typography, and
alignment rules.  Components use *python-pptx* primitives internally
and are called from individual layout renderers.

Components accept **grid spans** and **spacing tokens** rather than raw
coordinates.  They enforce text limits and hierarchy automatically.
"""

from __future__ import annotations

from pptx.util import Inches, Pt  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore

from ppt.generator import add_text_box, hex_to_rgb
from ppt.design_system import (
    ContentTransform,
    Grid,
    Spacing,
    Typography,
    VLayout,
    VStack,
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
)


# ======================================================================
# Accent / Decorative Elements
# ======================================================================

def accent_bar_top(slide, theme, *, height: float = 0.10):
    """Full-width accent bar at the very top."""
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme.primary)
    bar.line.fill.background()
    return bar


def accent_bar_bottom(slide, theme, *, height: float = 0.12,
                      color: str | None = None):
    """Full-width accent bar at the very bottom."""
    y = SLIDE_HEIGHT - height
    c = color or theme.accent
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(y), Inches(SLIDE_WIDTH), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(c)
    bar.line.fill.background()
    return bar


def accent_underline(slide, theme, *, left: float | None = None,
                     y: float | None = None,
                     width: float | None = None,
                     color: str | None = None):
    """Thin accent underline (typically below a heading).

    Width defaults to 2 grid columns for visual rhythm.
    """
    _left = left if left is not None else Grid.MARGIN
    _y = y if y is not None else VLayout.ACCENT_Y
    _w = width if width is not None else Grid.span_width(2)
    c = color or theme.accent
    bar = slide.shapes.add_shape(
        1, Inches(_left), Inches(_y), Inches(_w), Inches(VLayout.ACCENT_HEIGHT),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(c)
    bar.line.fill.background()
    return bar


def divider_line(slide, theme, *, y: float = 4.5,
                 span: int = 3, center: bool = True,
                 color: str | None = None):
    """Thin horizontal divider spanning *span* grid columns."""
    c = color or theme.accent
    width = Grid.span_width(span)
    left = (SLIDE_WIDTH - width) / 2 if center else Grid.MARGIN
    bar = slide.shapes.add_shape(
        1, Inches(left), Inches(y), Inches(width), Inches(0.05),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(c)
    bar.line.fill.background()
    return bar


def vertical_divider(slide, theme, *, y: float | None = None,
                     height: float | None = None,
                     color: str | None = None):
    """Vertical divider at the horizontal centre of the slide."""
    c = color or theme.accent
    _y = y if y is not None else VLayout.CONTENT_START
    _h = height if height is not None else (VLayout.CONTENT_END - _y)
    x = SLIDE_WIDTH / 2 - 0.03
    div = slide.shapes.add_shape(
        1, Inches(x), Inches(_y), Inches(0.06), Inches(_h),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = hex_to_rgb(c)
    div.line.fill.background()
    return div


# ======================================================================
# Text Components
# ======================================================================

def hero_text(slide, text: str, theme, *, y: float | None = None,
              span: int = 10, color: str | None = None,
              align=PP_ALIGN.CENTER):
    """Large hero text — the PRIMARY (dominant) visual element.

    Accepts *span* (grid columns) instead of raw width.
    Truncates text to enforce readability.
    """
    _y = y if y is not None else 2.4
    left, width = Grid.center(span)
    c = color or theme.primary
    display = ContentTransform.truncate(text, max_words=15)
    return add_text_box(
        slide, display,
        Inches(left), Inches(_y), Inches(width), Inches(1.6),
        theme.font_heading, Typography.HERO, bold=True,
        color=c, align=align,
    )


def heading_block(slide, text: str, theme, *,
                  span: int = 12, offset: int = 0,
                  y: float | None = None,
                  size: int | None = None,
                  color: str | None = None,
                  align=PP_ALIGN.LEFT):
    """Standard section heading at HEADING level.

    Accepts *span* and *offset* (grid columns).
    """
    left, width = Grid.compute(span, offset)
    _y = y if y is not None else VLayout.TITLE_TOP
    _sz = size or Typography.HEADING
    c = color or theme.primary
    display = ContentTransform.truncate(text, max_words=12)
    return add_text_box(
        slide, display,
        Inches(left), Inches(_y), Inches(width), Inches(VLayout.TITLE_HEIGHT),
        theme.font_heading, _sz, bold=True,
        color=c, align=align,
    )


def subheading_text(slide, text: str, theme, *,
                    span: int = 12, offset: int = 0,
                    y: float | None = None,
                    color: str | None = None,
                    align=PP_ALIGN.LEFT):
    """SECONDARY heading below the main heading."""
    left, width = Grid.compute(span, offset)
    _y = y if y is not None else VLayout.ACCENT_Y + Spacing.MD
    c = color or theme.secondary
    display = ContentTransform.truncate(text, max_words=15)
    return add_text_box(
        slide, display,
        Inches(left), Inches(y if y is not None else _y),
        Inches(width), Inches(0.7),
        theme.font_heading, Typography.SUBHEADING, bold=True,
        color=c, align=align,
    )


def body_text(slide, text: str, theme, *,
              left: float | None = None, y: float | None = None,
              width: float | None = None, height: float = 1.5,
              span: int | None = None, offset: int = 0,
              size: int | None = None,
              color: str | None = None, align=PP_ALIGN.LEFT):
    """Standard body text at BODY level.

    Accepts either explicit *left*/*width* OR *span*/*offset* (grid).
    """
    if span is not None:
        _left, _w = Grid.compute(span, offset)
    else:
        _left = left if left is not None else Grid.MARGIN
        _w = width if width is not None else Grid.USABLE_WIDTH
    _y = y if y is not None else VLayout.CONTENT_START
    _sz = size or Typography.BODY
    c = color or theme.text
    return add_text_box(
        slide, text,
        Inches(_left), Inches(_y), Inches(_w), Inches(height),
        theme.font_body, _sz, bold=False,
        color=c, align=align,
    )


def caption_text(slide, text: str, theme, *,
                 left: float | None = None, y: float | None = None,
                 width: float | None = None,
                 span: int | None = None,
                 color: str | None = None, align=PP_ALIGN.CENTER):
    """TERTIARY caption or source attribution."""
    if span is not None:
        _left, _w = Grid.center(span)
    else:
        _left = left if left is not None else Grid.MARGIN
        _w = width if width is not None else Grid.USABLE_WIDTH
    _y = y if y is not None else VLayout.CONTENT_END + Spacing.SM
    c = color or theme.secondary
    return add_text_box(
        slide, text,
        Inches(_left), Inches(_y), Inches(_w), Inches(0.45),
        theme.font_body, Typography.CAPTION, bold=False,
        color=c, align=align,
    )


# ======================================================================
# Bullet Group
# ======================================================================

def bullet_group(slide, items: list[str], theme, *,
                 left: float | None = None, start_y: float | None = None,
                 width: float | None = None,
                 span: int | None = None, offset: int = 0,
                 size: int | None = None, gap: float | None = None,
                 color: str | None = None,
                 max_items: int = 6, max_words: int = 12):
    """Render a truncated bullet list with consistent spacing."""
    if span is not None:
        _left, _w = Grid.compute(span, offset)
        _left += Spacing.SM
        _w -= Spacing.SM
    else:
        _left = left if left is not None else Grid.MARGIN + Spacing.MD
        _w = width if width is not None else Grid.USABLE_WIDTH - Spacing.MD
    _y = start_y if start_y is not None else VLayout.CONTENT_START
    _sz = size or Typography.BODY
    _gap = gap if gap is not None else Spacing.ELEMENT
    c = color or theme.text

    clean = ContentTransform.truncate_bullets(items, max_items=max_items,
                                              max_words=max_words)
    boxes = []
    for i, item in enumerate(clean):
        y_pos = _y + i * (0.6 + _gap)
        box = add_text_box(
            slide, f"•  {item}",
            Inches(_left), Inches(y_pos), Inches(_w), Inches(0.6),
            theme.font_body, _sz, bold=False,
            color=c, align=PP_ALIGN.LEFT,
        )
        boxes.append(box)
    return boxes


# ======================================================================
# Card
# ======================================================================

def card(slide, theme, *, left: float, y: float, width: float,
         height: float, icon: str = "", label: str = "",
         description: str = "", show_accent_strip: bool = False):
    """Card component with optional icon, label, and description.

    Text is auto-truncated for readability.
    """
    # Background
    bg = slide.shapes.add_shape(
        1, Inches(left), Inches(y), Inches(width), Inches(height),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(theme.surface)
    bg.line.color.rgb = hex_to_rgb(theme.secondary)
    bg.line.width = Pt(1)

    if show_accent_strip:
        strip = slide.shapes.add_shape(
            1, Inches(left), Inches(y), Inches(width), Inches(0.08),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = hex_to_rgb(theme.accent)
        strip.line.fill.background()

    pad = Spacing.SM
    inner_left = left + pad
    inner_w = width - 2 * pad
    cur_y = y + Spacing.MD

    if icon:
        add_text_box(
            slide, icon,
            Inches(left), Inches(cur_y), Inches(width), Inches(0.6),
            theme.font_body, Typography.SUBHEADING, bold=False,
            color=theme.accent, align=PP_ALIGN.CENTER,
        )
        cur_y += 0.6 + Spacing.SM

    if label:
        display_label = ContentTransform.truncate(label, max_words=6)
        add_text_box(
            slide, display_label,
            Inches(inner_left), Inches(cur_y), Inches(inner_w), Inches(0.6),
            theme.font_heading, Typography.BODY, bold=True,
            color=theme.primary, align=PP_ALIGN.CENTER,
        )
        cur_y += 0.6 + Spacing.SM

    if description:
        remaining = (y + height) - cur_y - pad
        display_desc = ContentTransform.truncate(description, max_words=20)
        add_text_box(
            slide, display_desc,
            Inches(inner_left), Inches(cur_y),
            Inches(inner_w), Inches(max(remaining, 0.5)),
            theme.font_body, Typography.CAPTION, bold=False,
            color=theme.text, align=PP_ALIGN.CENTER,
        )


# ======================================================================
# Stat Block
# ======================================================================

def stat_block(slide, stat: str, label: str, theme, *,
               y: float | None = None, span: int = 10,
               left: float | None = None, width: float | None = None):
    """Large statistic — PRIMARY dominant element.

    Accepts *span* (grid columns) for width.
    """
    _y = y if y is not None else VLayout.CONTENT_START
    if left is not None and width is not None:
        _left, _width = left, width
    else:
        _left, _width = Grid.center(span)

    # Big number
    add_text_box(
        slide, stat,
        Inches(_left), Inches(_y), Inches(_width), Inches(1.8),
        theme.font_heading, Typography.HERO + 12, bold=True,
        color=theme.accent, align=PP_ALIGN.CENTER,
    )
    # Label underneath
    display_label = ContentTransform.truncate(label, max_words=8)
    add_text_box(
        slide, display_label,
        Inches(_left), Inches(_y + 1.8 + Spacing.SM),
        Inches(_width), Inches(0.7),
        theme.font_body, Typography.SUBHEADING - 4, bold=False,
        color=theme.secondary, align=PP_ALIGN.CENTER,
    )


# ======================================================================
# Two-Column Helpers
# ======================================================================

def two_column_headers(slide, left_label: str, right_label: str, theme,
                       *, y: float | None = None):
    """Render column headers using Grid.split(6, 6)."""
    _y = y if y is not None else VLayout.CONTENT_START
    cols = Grid.split(6, 6)

    for i, label_text in enumerate((left_label, right_label)):
        display = ContentTransform.truncate(label_text, max_words=6)
        add_text_box(
            slide, display,
            Inches(cols[i][0]), Inches(_y),
            Inches(cols[i][1]), Inches(0.7),
            theme.font_heading, Typography.SUBHEADING - 4, bold=True,
            color=theme.primary, align=PP_ALIGN.CENTER,
        )
    return cols


# ======================================================================
# Highlight Box
# ======================================================================

def highlight_box(slide, theme, *, label: str, text: str,
                  y: float | None = None, height: float = 1.4):
    """Coloured highlight box (e.g. key takeaway) at full grid width."""
    _y = y if y is not None else VLayout.CONTENT_END - height
    left, width = Grid.full_width()
    bg = slide.shapes.add_shape(
        1, Inches(left), Inches(_y), Inches(width), Inches(height),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(theme.primary)
    bg.line.fill.background()

    inner_left = left + Spacing.MD
    inner_w = width - 2 * Spacing.MD

    add_text_box(
        slide, label,
        Inches(inner_left), Inches(_y + Spacing.SM),
        Inches(inner_w), Inches(0.4),
        theme.font_heading, Typography.CAPTION, bold=True,
        color=theme.background, align=PP_ALIGN.LEFT,
    )
    display = ContentTransform.truncate(text, max_words=25)
    add_text_box(
        slide, display,
        Inches(inner_left), Inches(_y + Spacing.SM + 0.45),
        Inches(inner_w), Inches(height - 0.65),
        theme.font_body, Typography.BODY, bold=False,
        color=theme.background, align=PP_ALIGN.LEFT,
    )
