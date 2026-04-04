"""
Stage 4: Export Engine

Supports BOTH modern and legacy export paths:

Modern:
  - HTML slides (standalone per slide)
  - PNG images (from rendering engine)
  - PDF (combined document)

Legacy (python-pptx):
  - Each slide rendered as an image → inserted into PPT as image
  - PPT is only a container for images (no shapes/text)
  - Falls back to the standard PPT generator if no images available
"""

import logging
import os
import uuid
from typing import Optional

logger = logging.getLogger(__name__)


def build_ppt_structure(image_paths: list[str]) -> dict:
    """
    Build the PPT_STRUCTURE output for legacy export.

    Each slide becomes an image reference.
    """
    slides = []
    for path in image_paths:
        slides.append({"image": os.path.basename(path)})

    return {"slides": slides}


def generate_image_based_ppt(
    image_paths: list[str],
    output_dir: str,
) -> Optional[str]:
    """
    Generate a PPT file where each slide contains a full-bleed image.

    This is the legacy export path — PPT is only a container for
    rendered slide images. No design via pptx shapes.

    Returns the PPT file path, or None if no images or pptx unavailable.
    """
    if not image_paths:
        logger.info("[export_engine] No images available — skipping image-based PPT")
        return None

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Emu  # type: ignore
    except ImportError:
        logger.warning("[export_engine] python-pptx not installed — skipping PPT export")
        return None

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for img_path in image_paths:
        if not os.path.isfile(img_path):
            logger.warning("[export_engine] Image not found: %s", img_path)
            continue

        slide = prs.slides.add_slide(blank_layout)
        # Full-bleed image covering entire slide
        slide.shapes.add_picture(
            img_path,
            Emu(0),
            Emu(0),
            prs.slide_width,
            prs.slide_height,
        )

    os.makedirs(output_dir, exist_ok=True)
    ppt_path = os.path.join(output_dir, f"{uuid.uuid4().hex}_visual.pptx")
    prs.save(ppt_path)
    logger.info("[export_engine] Generated image-based PPT: %s", ppt_path)
    return ppt_path


def save_html_slides(
    html_slides: list[str],
    output_dir: str,
) -> list[str]:
    """Save individual HTML slides to disk."""
    os.makedirs(output_dir, exist_ok=True)
    paths = []
    for idx, slide_html in enumerate(html_slides):
        path = os.path.join(output_dir, f"slide_{idx + 1}.html")
        with open(path, "w", encoding="utf-8") as f:
            f.write(slide_html)
        paths.append(path)
    logger.info("[export_engine] Saved %d HTML slides", len(paths))
    return paths


def run_export_engine(
    html_slides: list[str],
    image_paths: list[str],
    pdf_path: Optional[str],
    output_dir: str,
) -> dict:
    """
    Stage 4 entry point.

    Produces all export artifacts:
      - HTML slide files
      - PPT structure (for legacy)
      - Image-based PPT file (if images available)

    Returns:
        {
            "html_paths": [...],
            "image_paths": [...],
            "pdf_path": str | None,
            "ppt_path": str | None,
            "ppt_structure": { "slides": [...] },
        }
    """
    # Modern: save HTML
    html_paths = save_html_slides(html_slides, output_dir)

    # Legacy: build PPT structure and generate image-based PPT
    ppt_structure = build_ppt_structure(image_paths)
    ppt_path = generate_image_based_ppt(image_paths, output_dir)

    return {
        "html_paths": html_paths,
        "image_paths": image_paths,
        "pdf_path": pdf_path,
        "ppt_path": ppt_path,
        "ppt_structure": ppt_structure,
    }
