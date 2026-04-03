"""Tests for PPT generation and themes."""
import os
import pytest
from pptx import Presentation

from models.presentation_state import PresentationState
from pipeline.state_builder import build_state
from pipeline.slide_planner import plan_slides
from pipeline.slide_type_assigner import assign_slide_types
from pipeline.story_generator import _default_story
from ppt.generator import generate_ppt, hex_to_rgb, add_text_box, set_background
from ppt.themes.modern import THEMES, ThemeConfig


# --- Mock content generator ---
MOCK_CONTENT = {
    "title_slide": {"title": "Test Presentation", "subtitle": "Subtitle", "presenter": "Tester"},
    "section_header": {"section_title": "Section", "tagline": "Tag"},
    "agenda_slide": {"title": "Agenda", "items": ["Item 1", "Item 2", "Item 3"]},
    "problem_slide": {"title": "Problem", "cards": [{"icon": "⚠️", "label": "Issue", "description": "Desc"}]},
    "stats_slide": {"title": "Stats", "stat": "42%", "stat_label": "Label", "description": "Desc", "source": "Source"},
    "feature_slide": {"title": "Features", "features": [{"icon": "🔬", "label": "Feature", "description": "Desc"}]},
    "comparison_slide": {"title": "Compare", "left_label": "A", "left_points": ["a1"], "right_label": "B", "right_points": ["b1"]},
    "timeline_slide": {"title": "Timeline", "events": [{"year": "2024", "label": "Event"}]},
    "example_slide": {"title": "Example", "example_title": "Case", "context": "Ctx", "result": "Res", "takeaway": "Key"},
    "quote_slide": {"quote": "Innovation matters.", "attribution": "Author"},
    "image_slide": {"title": "Image", "caption": "Caption"},
    "conclusion_slide": {"title": "Conclusion", "bullets": ["Point 1"], "key_takeaway": "Takeaway"},
    "cta_slide": {"title": "Act Now", "cta_text": "Contact us", "contact": "test@test.com"},
    "thank_you_slide": {"title": "Thank You", "message": "Questions?", "contact": "test@test.com"},
}


def _build_test_state(theme="modern"):
    state = build_state({
        "topic": "Test Topic",
        "presentation_type": "pitch",
        "slide_count": 8,
        "sections": ["intro", "problem", "solution", "benefits", "conclusion"],
        "tone": "professional",
        "theme": theme,
    })
    state = state.model_copy(update={"story": _default_story(state)})
    state = plan_slides(state)
    state = assign_slide_types(state)

    structured = []
    for slide in state.slide_plan:
        t = slide["type"]
        content = MOCK_CONTENT.get(t, {"title": "Fallback", "body": "Content"})
        structured.append({
            "slide_id": slide["slide_id"],
            "type": t,
            "content": content,
            "image_path": None,
        })
    state = state.model_copy(update={"structured_slides": structured})

    notes = [{"slide_id": s["slide_id"], "notes": f"Notes for slide {s['slide_id']}"} for s in structured]
    state = state.model_copy(update={"speaker_notes": notes})
    return state


def test_generate_ppt_creates_file():
    state = _build_test_state()
    path = generate_ppt(state)
    assert os.path.isfile(path)
    assert path.endswith(".pptx")
    os.remove(path)


def test_generate_ppt_correct_slide_count():
    state = _build_test_state()
    path = generate_ppt(state)
    prs = Presentation(path)
    assert len(prs.slides) == len(state.structured_slides)
    os.remove(path)


def test_generate_ppt_speaker_notes_injected():
    state = _build_test_state()
    path = generate_ppt(state)
    prs = Presentation(path)
    notes_found = 0
    for slide in prs.slides:
        try:
            ns = slide.notes_slide
            if ns and ns.notes_text_frame.text:
                notes_found += 1
        except Exception:
            pass
    assert notes_found == len(state.structured_slides)
    os.remove(path)


def test_generate_ppt_widescreen():
    state = _build_test_state()
    path = generate_ppt(state)
    prs = Presentation(path)
    # 16:9 widescreen (13.33" x 7.5")
    assert prs.slide_width > prs.slide_height
    os.remove(path)


@pytest.mark.parametrize("theme_name", ["modern", "corporate", "minimal"])
def test_generate_ppt_all_themes(theme_name):
    state = _build_test_state(theme=theme_name)
    path = generate_ppt(state)
    assert os.path.isfile(path)
    prs = Presentation(path)
    assert len(prs.slides) > 0
    os.remove(path)


def test_hex_to_rgb():
    rgb = hex_to_rgb("6C63FF")
    assert rgb.red == 0x6C
    assert rgb.green == 0x63
    assert rgb.blue == 0xFF


def test_hex_to_rgb_with_hash():
    rgb = hex_to_rgb("#FF0000")
    assert rgb.red == 0xFF
    assert rgb.green == 0x00
    assert rgb.blue == 0x00


def test_all_themes_exist():
    assert "modern" in THEMES
    assert "corporate" in THEMES
    assert "minimal" in THEMES


def test_theme_config_fields():
    for name, theme in THEMES.items():
        assert isinstance(theme, ThemeConfig)
        assert theme.name == name
        assert theme.primary
        assert theme.font_heading
        assert theme.font_body
        assert theme.heading_size > 0
        assert theme.body_size > 0
