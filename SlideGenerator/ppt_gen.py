import os
from pptx import Presentation # type: ignore

def create_ppt(slides_data, file_path):
    # Ensure directory exists
    os.makedirs(os.path.dirname(file_path), exist_ok=True)

    prs = Presentation()

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data["title"]
        content.text = "\n".join(slide_data["bullets"])

    prs.save(file_path)